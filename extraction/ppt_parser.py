from pptx import Presentation
import os


def extract_text_from_ppt(file_path):

    slides = []

    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return []

    if not file_path.lower().endswith(".pptx"):
        print(f"Unsupported format: only .pptx is supported, got: {file_path}")
        return []

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("Failed to open PPTX:", e)
        return []

    for i, slide in enumerate(prs.slides):

        slide_parts = []
        slide_title = ""

        for shape in slide.shapes:

            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                is_title = False
                if shape.is_placeholder:
                    try:
                        is_title = shape.placeholder_format.idx == 0
                    except Exception:
                        pass
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        if is_title and not slide_title:
                            slide_title = line
                        else:
                            slide_parts.append(line)

            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    seen = set()
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and cell_text not in seen:
                            seen.add(cell_text)
                            slide_parts.append(cell_text)

        # Normalize whitespace
        slide_text = " ".join(slide_parts)
        slide_text = " ".join(slide_text.split())

        if slide_text or slide_title:
            slides.append({
                "slide": i + 1,
                "heading": slide_title,
                "text": slide_text
            })

    return slides


if __name__ == "__main__":

    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

    sample_file = os.path.join(base_dir, "data", "uploads", "sample.pptx")

    slides = extract_text_from_ppt(sample_file)

    full_text = " ".join(s["text"] for s in slides)
    print(full_text[:1000])